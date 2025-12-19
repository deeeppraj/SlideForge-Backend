import os
import uuid
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(*tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4)))


def clamp_text(text, max_chars):
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rsplit(" ", 1)[0] + "…"


class ColorPalette:
    BACKGROUND = "0A0E27"
    CARD_BG = "0F1729"
    ACCENT_CYAN = "00D9FF"
    ACCENT_PURPLE = "9D4EDD"
    ACCENT_PINK = "FF006E"
    TEXT_PRIMARY = "F5F7FA"
    TEXT_SECONDARY = "B8BEC9"
    TEXT_MUTED = "7C8696"
    BORDER_COLOR = "1B2845"


class TypographyConfig:
    FONT_PRIMARY = "Calibri"
    FONT_SIZE_TITLE = Pt(54)
    FONT_SIZE_POINT_4 = Pt(22)
    FONT_SIZE_POINT_3 = Pt(28)
    FONT_SIZE_POINT_2 = Pt(32)
    FONT_SIZE_POINT_1 = Pt(36)
    FONT_SIZE_EXP_4 = Pt(12)
    FONT_SIZE_EXP_3 = Pt(14)
    FONT_SIZE_EXP_2 = Pt(15)
    FONT_SIZE_EXP_1 = Pt(16)


def create_elegant_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(ColorPalette.BACKGROUND)
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    
    glow_top_left = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-1.2), Inches(-0.8),
        Inches(4.5), Inches(4.5)
    )
    glow_top_left.fill.solid()
    glow_top_left.fill.fore_color.rgb = hex_to_rgb(ColorPalette.ACCENT_CYAN)
    glow_top_left.fill.transparency = 0.93
    glow_top_left.line.fill.background()
    
    glow_bottom_right = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(8.8), Inches(4.2),
        Inches(4.2), Inches(4.2)
    )
    glow_bottom_right.fill.solid()
    glow_bottom_right.fill.fore_color.rgb = hex_to_rgb(ColorPalette.ACCENT_PURPLE)
    glow_bottom_right.fill.transparency = 0.92
    glow_bottom_right.line.fill.background()
    
    glow_accent = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4.8), Inches(2.5),
        Inches(2), Inches(2)
    )
    glow_accent.fill.solid()
    glow_accent.fill.fore_color.rgb = hex_to_rgb(ColorPalette.ACCENT_PINK)
    glow_accent.fill.transparency = 0.94
    glow_accent.line.fill.background()
    
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(0.2),
        Inches(9.4), Inches(1.5)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = hex_to_rgb(ColorPalette.CARD_BG)
    title_bg.fill.transparency = 0.3
    title_bg.line.fill.background()
    
    title_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(1.65),
        Inches(9.4), Inches(0.04)
    )
    title_accent.fill.solid()
    title_accent.fill.fore_color.rgb = hex_to_rgb(ColorPalette.ACCENT_CYAN)
    title_accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
    Inches(0.5), Inches(1.8), Inches(9), Inches(2.2)  
    )
    tf = title_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    title = tf.paragraphs[0]
    title.text = clamp_text(slide_data["title"], 85)
    title.font.size = TypographyConfig.FONT_SIZE_TITLE
    title.font.bold = True
    title.font.color.rgb = hex_to_rgb(ColorPalette.TEXT_PRIMARY)
    title.font.name = TypographyConfig.FONT_PRIMARY
    
    content_top = 2.0
    content_height = 5.0
    
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(content_top),
        Inches(6.2), Inches(content_height)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = hex_to_rgb(ColorPalette.CARD_BG)
    left_panel.fill.transparency = 0.15
    left_panel.line.color.rgb = hex_to_rgb(ColorPalette.BORDER_COLOR)
    left_panel.line.width = Pt(0.75)
    left_panel.adjustments[0] = 0.08
    
    text_box = slide.shapes.add_textbox(
        Inches(0.55), Inches(content_top + 0.25),
        Inches(5.7), Inches(content_height - 0.5)
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()
    
    points = slide_data.get("points", [])
    explanations = slide_data.get("explanation", [])
    count = min(len(points), len(explanations), 4)
    
    config_map = {
        1: {
            "point_size": TypographyConfig.FONT_SIZE_POINT_1,
            "exp_size": TypographyConfig.FONT_SIZE_EXP_1,
            "max_chars": 150,
            "line_spacing": 1.3,
            "vertical_spacing": 0.15
        },
        2: {
            "point_size": TypographyConfig.FONT_SIZE_POINT_2,
            "exp_size": TypographyConfig.FONT_SIZE_EXP_2,
            "max_chars": 130,
            "line_spacing": 1.3,
            "vertical_spacing": 0.25
        },
        3: {
            "point_size": TypographyConfig.FONT_SIZE_POINT_3,
            "exp_size": TypographyConfig.FONT_SIZE_EXP_3,
            "max_chars": 115,
            "line_spacing": 1.4,
            "vertical_spacing": 0.35
        },
        4: {
                "point_size": Pt(22),       
                "exp_size": Pt(12),          
                "line_spacing": 1.25,        
                "space_before": 0,          
                "space_after_point": 1,      
                "space_after_exp": 3      
            }
    }
    
    cfg = config_map[count]
    
    for i in range(count):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "▸ " + clamp_text(points[i], 55)
        p.font.size = cfg["point_size"]
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(ColorPalette.ACCENT_CYAN)
        p.font.name = TypographyConfig.FONT_PRIMARY
        p.space_before = Pt(cfg["vertical_spacing"] * 72) if i > 0 else Pt(0)
        p.space_after = Pt(6)
        
        exp = tf.add_paragraph()
        exp.text = clamp_text(explanations[i], cfg["max_chars"])
        exp.font.size = cfg["exp_size"]
        exp.font.color.rgb = hex_to_rgb(ColorPalette.TEXT_SECONDARY)
        exp.font.name = TypographyConfig.FONT_PRIMARY
        exp.line_spacing = cfg["line_spacing"]
        exp.space_after = Pt(0)
    
    img_size = Inches(2.9)
    img_left = Inches(6.7)
    img_top = Inches(content_top + 0.3)
    
    img_frame = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        img_left, img_top, img_size, img_size
    )
    img_frame.fill.solid()
    img_frame.fill.fore_color.rgb = hex_to_rgb(ColorPalette.CARD_BG)
    img_frame.fill.transparency = 0.05
    img_frame.line.color.rgb = hex_to_rgb(ColorPalette.ACCENT_PURPLE)
    img_frame.line.width = Pt(1.5)
    img_frame.adjustments[0] = 0.08
    
    try:
        if "image" in slide_data and slide_data["image"]:
            img_bytes = requests.get(slide_data["image"], timeout=5).content
            path = f"{uuid.uuid4()}.jpg"
            with open(path, "wb") as f:
                f.write(img_bytes)
            
            slide.shapes.add_picture(
                path,
                img_left + Inches(0.12),
                img_top + Inches(0.12),
                width=img_size - Inches(0.24),
                height=img_size - Inches(0.24)
            )
            
            os.remove(path)
    except Exception as e:
        pass
    
    footer_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.3), Inches(7.35),
        Inches(9.4), Inches(0.03)
    )
    footer_line.fill.solid()
    footer_line.fill.fore_color.rgb = hex_to_rgb(ColorPalette.ACCENT_CYAN)
    footer_line.line.fill.background()


def create_title_slide(prs, title, subtitle="", author=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(ColorPalette.BACKGROUND)
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    
    glow1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-1.5), Inches(-1),
        Inches(5), Inches(5)
    )
    glow1.fill.solid()
    glow1.fill.fore_color.rgb = hex_to_rgb(ColorPalette.ACCENT_CYAN)
    glow1.fill.transparency = 0.92
    glow1.line.fill.background()
    
    glow2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(8.5), Inches(4),
        Inches(4.5), Inches(4.5)
    )
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = hex_to_rgb(ColorPalette.ACCENT_PURPLE)
    glow2.fill.transparency = 0.91
    glow2.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clamp_text(title, 100)
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(ColorPalette.TEXT_PRIMARY)
    p.font.name = TypographyConfig.FONT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.2), Inches(8), Inches(1)
        )
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = clamp_text(subtitle, 120)
        p.font.size = Pt(28)
        p.font.color.rgb = hex_to_rgb(ColorPalette.ACCENT_CYAN)
        p.font.name = TypographyConfig.FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
    
    if author:
        author_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(8), Inches(0.5)
        )
        tf = author_box.text_frame
        p = tf.paragraphs[0]
        p.text = author
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(ColorPalette.TEXT_MUTED)
        p.font.name = TypographyConfig.FONT_PRIMARY
        p.alignment = PP_ALIGN.CENTER

